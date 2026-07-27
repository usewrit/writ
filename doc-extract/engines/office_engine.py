"""Office documents — docx / xlsx / pptx. Pure text-layer extraction, never OCR.

These formats always carry their text as data, so we parse them directly with
the format-native libraries and emit markdown + structured records. OCR would be
strictly worse and is never used here.
"""
from __future__ import annotations

import io
import logging
from typing import List

logger = logging.getLogger("doc-extract.office")


def _md_table(rows: List[list]) -> str:
    """Render a list-of-rows as a GitHub markdown table (best-effort)."""
    rows = [[("" if c is None else str(c)) for c in r] for r in rows if r is not None]
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    head = rows[0]
    body = rows[1:] if len(rows) > 1 else []
    lines = ["| " + " | ".join(head) + " |", "| " + " | ".join(["---"] * width) + " |"]
    for r in body:
        lines.append("| " + " | ".join(r) + " |")
    return "\n".join(lines)


def _extract_docx(data: bytes, source_url: str) -> dict:
    import docx  # python-docx

    document = docx.Document(io.BytesIO(data))
    md_parts: List[str] = []
    text_parts: List[str] = []
    for para in document.paragraphs:
        t = (para.text or "").strip()
        if not t:
            continue
        style = (para.style.name or "").lower() if para.style else ""
        if style.startswith("heading"):
            level = "".join(ch for ch in style if ch.isdigit()) or "2"
            md_parts.append(f"{'#' * min(int(level), 6)} {t}")
        else:
            md_parts.append(t)
        text_parts.append(t)

    tables: List[list] = []
    for table in document.tables:
        rows = [[(cell.text or "").strip() for cell in row.cells] for row in table.rows]
        if rows:
            tables.append(rows)
            md_parts.append(_md_table(rows))

    return {
        "kind": "docx",
        "content_kind": "docx",
        "markdown": "\n\n".join(md_parts).strip(),
        "text": "\n".join(text_parts).strip(),
        "tables": tables,
        "records": [],
        "ocr": None,
        "meta": {"source_url": source_url},
    }


def _extract_xlsx(data: bytes, source_url: str) -> dict:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    md_parts: List[str] = []
    tables: List[list] = []
    records: List[dict] = []
    for ws in wb.worksheets:
        rows = [
            [("" if c is None else c) for c in row]
            for row in ws.iter_rows(values_only=True)
        ]
        rows = [r for r in rows if any(str(c).strip() for c in r)]
        if not rows:
            continue
        tables.append(rows)
        md_parts.append(f"## {ws.title}\n\n" + _md_table(rows))
        header = [str(c).strip() for c in rows[0]]
        for r in rows[1:]:
            records.append({header[i] if i < len(header) else f"col{i}": r[i]
                            for i in range(len(r))})
    wb.close()
    return {
        "kind": "xlsx",
        "content_kind": "xlsx",
        "markdown": "\n\n".join(md_parts).strip(),
        "text": "\n\n".join(md_parts).strip(),
        "tables": tables,
        "records": records,
        "ocr": None,
        "meta": {"source_url": source_url, "sheets": len(tables)},
    }


def _extract_pptx(data: bytes, source_url: str) -> dict:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    md_parts: List[str] = []
    text_parts: List[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        md_parts.append(f"## Slide {idx}")
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                t = "".join(run.text for run in para.runs).strip()
                if t:
                    md_parts.append(t)
                    text_parts.append(t)
    return {
        "kind": "pptx",
        "content_kind": "pptx",
        "markdown": "\n\n".join(md_parts).strip(),
        "text": "\n".join(text_parts).strip(),
        "tables": [],
        "records": [],
        "ocr": None,
        "meta": {"source_url": source_url, "slides": len(prs.slides._sldIdLst)},
    }


def extract(data: bytes, source_url: str, kind: str) -> dict:
    if kind == "docx":
        return _extract_docx(data, source_url)
    if kind == "xlsx":
        return _extract_xlsx(data, source_url)
    if kind == "pptx":
        return _extract_pptx(data, source_url)
    raise ValueError(f"unsupported office kind: {kind}")
